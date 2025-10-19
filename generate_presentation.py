"""
Generate a PowerPoint deck from Netflix analysis
- Loads cleaned or raw dataset
- Creates key charts as PNGs
- Builds a PPTX with executive summary and visuals
"""

import os
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
from datetime import datetime

# pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

sns.set(style="whitegrid")

DATA_PATHS = [
    'outputs/cleaned_netflix.csv',
    'Netflix Dataset.csv',
]

OUT_DIR = 'outputs'
CHART_DIR = os.path.join(OUT_DIR, 'charts')
os.makedirs(OUT_DIR, exist_ok=True)
os.makedirs(CHART_DIR, exist_ok=True)


def load_data():
    path = None
    for p in DATA_PATHS:
        if os.path.exists(p):
            path = p
            break
    if not path:
        raise FileNotFoundError(f"Dataset not found. Checked {DATA_PATHS}")
    df = pd.read_csv(path)
    df.columns = [c.strip().lower().replace(' ', '_') for c in df.columns]
    # Map columns if raw csv
    if 'category' in df.columns:
        df = df.rename(columns={'category': 'type'})
    if 'type' in df.columns and 'listed_in' not in df.columns:
        df = df.rename(columns={'type': 'listed_in'})
    if 'release_date' in df.columns:
        df = df.rename(columns={'release_date': 'date_added'})
    df['date_added'] = pd.to_datetime(df.get('date_added'), errors='coerce')
    df['year_added'] = df['date_added'].dt.year
    return df


def save_matplotlib(fig, name):
    path = os.path.join(CHART_DIR, name)
    fig.savefig(path, bbox_inches='tight', dpi=200)
    plt.close(fig)
    return path

def save_plotly(fig, name):
    """Save plotly figure as PNG using kaleido if available; return path or None."""
    path = os.path.join(CHART_DIR, name)
    try:
        # Requires kaleido
        fig.write_image(path, scale=2)
        return path
    except Exception as e:
        print(f"[WARN] Could not export plotly image '{name}': {e}. Install 'kaleido' to enable static image export.")
        return None

def safe_mode(series, default='N/A'):
    try:
        m = series.mode(dropna=True)
        return m.iloc[0] if not m.empty else default
    except Exception:
        return default


def create_charts(df):
    charts = {}

    # 1) Movies vs TV Shows
    if 'type' in df.columns:
        vc = df['type'].value_counts()
        fig, ax = plt.subplots(figsize=(5, 3))
        sns.barplot(x=vc.index, y=vc.values, ax=ax, palette=['#E50914', '#564d4d'])
        ax.set_title('Distribution: Movies vs TV Shows')
        ax.set_xlabel('Type')
        ax.set_ylabel('Count')
        charts['type_dist'] = save_matplotlib(fig, 'type_distribution.png')

    # 2) Content added per year
    if 'year_added' in df.columns and 'type' in df.columns:
        year_counts = df.dropna(subset=['year_added']).groupby(['year_added','type']).size().reset_index(name='count')
        fig = px.area(year_counts, x='year_added', y='count', color='type', title='Content Added Per Year by Type')
        fig.update_layout(template='plotly_white')
        area_path = save_plotly(fig, 'content_per_year.png')
        if area_path:
            charts['per_year'] = area_path

    # 3) Top genres
    if 'listed_in' in df.columns:
        genres = df['listed_in'].dropna().str.split(', ').explode().value_counts().head(15)
        fig, ax = plt.subplots(figsize=(6, 4))
        sns.barplot(y=genres.index, x=genres.values, ax=ax, palette='Reds_r')
        ax.set_title('Top 15 Genres')
        ax.set_xlabel('Count')
        ax.set_ylabel('Genre')
        charts['top_genres'] = save_matplotlib(fig, 'top_genres.png')

    # 4) Top countries
    if 'country' in df.columns:
        countries = df['country'].dropna().str.split(', ').explode().value_counts().head(15)
        fig, ax = plt.subplots(figsize=(6, 4))
        sns.barplot(y=countries.index, x=countries.values, ax=ax, palette='Reds_r')
        ax.set_title('Top 15 Countries by Titles')
        ax.set_xlabel('Count')
        ax.set_ylabel('Country')
        charts['top_countries'] = save_matplotlib(fig, 'top_countries.png')

    # 5) Movie duration dist
    if 'duration' in df.columns and 'type' in df.columns:
        movies = df[df['type']=='Movie'].copy()
        if not movies.empty:
            movies['duration_minutes'] = movies['duration'].str.extract(r'(\d+)').astype(float)
            fig, ax = plt.subplots(figsize=(6, 3))
            ax.hist(movies['duration_minutes'].dropna(), bins=30, color='#E50914')
            ax.set_title('Distribution of Movie Durations (min)')
            ax.set_xlabel('Minutes')
            ax.set_ylabel('Count')
            charts['duration'] = save_matplotlib(fig, 'movie_duration.png')

    return charts


def slide_title(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def slide_bullets(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
    return slide


def slide_image(prs, title, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    left = Inches(1.0)
    top = Inches(1.7)
    width = Inches(8.0)
    slide.shapes.add_picture(img_path, left, top, width=width)
    return slide


def build_ppt(df, charts):
    prs = Presentation()

    # Title Slide
    slide_title(prs, 'Netflix Dataset Analysis', f"Auto-generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}\nSource: Netflix Dataset (7,789 rows)")

    # Executive Summary
    movies = int((df['type']=='Movie').sum()) if 'type' in df.columns else 0
    tvshows = int((df['type']=='TV Show').sum()) if 'type' in df.columns else 0
    top_genre = safe_mode(df['listed_in'].dropna().str.split(', ').explode()) if 'listed_in' in df.columns else 'N/A'
    top_country = safe_mode(df['country'].dropna().str.split(', ').explode()) if 'country' in df.columns else 'N/A'

    slide_bullets(
        prs,
        'Executive Summary',
        [
            f"Total titles: {len(df):,}",
            f"Movies vs TV Shows: {movies:,} vs {tvshows:,}",
            f"Top Genre: {top_genre}",
            f"Top Content Country: {top_country}",
            "Growth: Strong increase in content additions in recent years",
            "Recommendation: Invest in top genres and expand regional originals",
        ]
    )

    # Charts
    if 'type_dist' in charts:
        slide_image(prs, 'Movies vs TV Shows', charts['type_dist'])
    if 'per_year' in charts:
        slide_image(prs, 'Content Added Per Year by Type', charts['per_year'])
    if 'top_genres' in charts:
        slide_image(prs, 'Top Genres', charts['top_genres'])
    if 'top_countries' in charts:
        slide_image(prs, 'Top Countries by Titles', charts['top_countries'])
    if 'duration' in charts:
        slide_image(prs, 'Distribution of Movie Durations', charts['duration'])

    # Recommendations slide
    slide_bullets(
        prs,
        'Strategic Recommendations',
        [
            'Increase investment in TV Shows to boost retention',
            'Prioritize top-performing genres (Drama, Comedy, Documentary)',
            'Expand regional originals in underrepresented markets',
            'Maintain movie durations around 90–120 minutes',
            'Monitor quarterly trends; adjust acquisition strategy accordingly',
        ]
    )

    out_path = os.path.join(OUT_DIR, 'Netflix_Analysis_Presentation.pptx')
    prs.save(out_path)
    return out_path


def main():
    df = load_data()
    charts = create_charts(df)
    ppt_path = build_ppt(df, charts)
    print(f"Presentation created: {ppt_path}")


if __name__ == '__main__':
    main()
